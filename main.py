# pip3 install python-pptx
from pptx import Presentation

def extract_text_from_pptx(pptx_file):
    prs = Presentation(pptx_file)
    text = ""

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text += run.text

    return text